import os
import re
import pymupdf as fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.db import get_conn, DB_PATH

IMAGES_ROOT = os.path.join(os.path.dirname(DB_PATH), "images")
PDF_DPI = 140


def _lecture_dir(lecture_id):
    d = os.path.join(IMAGES_ROOT, str(lecture_id))
    os.makedirs(d, exist_ok=True)
    return d


def _rel_path(abs_path):
    """Path relative to IMAGES_ROOT, so /images/ + rel resolves correctly."""
    return os.path.relpath(abs_path, IMAGES_ROOT)


def _clear_slide_images(conn, lecture_id):
    conn.execute(
        "DELETE FROM slide_images WHERE slide_id IN "
        "(SELECT id FROM slides WHERE lecture_id=?)",
        (lecture_id,),
    )
    # remove orphan files
    d = _lecture_dir(lecture_id)
    if os.path.isdir(d):
        for f in os.listdir(d):
            try:
                os.remove(os.path.join(d, f))
            except OSError:
                pass


def _insert(conn, slide_id, abs_path, kind, seq):
    conn.execute(
        "INSERT INTO slide_images(slide_id, path, kind, seq) VALUES(?,?,?,?)",
        (slide_id, _rel_path(abs_path), kind, seq),
    )


def extract_pdf_images(lecture_id, path, conn=None):
    """Render each PDF page to a full-page JPEG."""
    own_conn = conn is None
    conn = conn or get_conn()
    try:
        slides = conn.execute(
            "SELECT id, slide_num FROM slides WHERE lecture_id=? ORDER BY slide_num",
            (lecture_id,),
        ).fetchall()
        slide_by_num = {r["slide_num"]: r["id"] for r in slides}

        doc = fitz.open(path)
        count = 0
        for pno in range(len(doc)):
            page = doc[pno]
            num = pno + 1
            if num not in slide_by_num:
                continue
            pix = page.get_pixmap(dpi=PDF_DPI)
            dest = os.path.join(_lecture_dir(lecture_id), f"page_{num}.jpg")
            pix.save(dest, jpg_quality=85)
            _insert(conn, slide_by_num[num], dest, "page", 0)
            count += 1
            # Commit per page so the write lock is never held across the slow loop.
            conn.commit()
        doc.close()
        if own_conn:
            conn.commit()
        return count
    finally:
        if own_conn:
            conn.close()


def _iter_picture_shapes(shapes):
    """Yield picture shapes, recursing into groups (best-effort)."""
    for shape in shapes:
        st = shape.shape_type
        if st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
            yield shape
        elif st == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            yield from _iter_picture_shapes(shape.shapes)


def extract_pptx_images(lecture_id, path, conn=None):
    """Extract embedded images per slide from a PPTX."""
    own_conn = conn is None
    conn = conn or get_conn()
    try:
        slides = conn.execute(
            "SELECT id, slide_num FROM slides WHERE lecture_id=? ORDER BY slide_num",
            (lecture_id,),
        ).fetchall()
        slide_by_num = {r["slide_num"]: r["id"] for r in slides}

        prs = Presentation(path)
        count = 0
        for idx, slide in enumerate(prs.slides):
            num = idx + 1
            if num not in slide_by_num:
                continue
            seq = 0
            for shape in _iter_picture_shapes(slide.shapes):
                try:
                    blob = shape.image.blob
                    ext = shape.image.ext or "png"
                except Exception:
                    continue
                dest = os.path.join(_lecture_dir(lecture_id), f"slide_{num}_{seq}.{ext}")
                with open(dest, "wb") as f:
                    f.write(blob)
                _insert(conn, slide_by_num[num], dest, "embedded", seq)
                seq += 1
                count += 1
                # Commit per image so the write lock is never held across the slow loop.
                conn.commit()
        if own_conn:
            conn.commit()
        return count
    finally:
        if own_conn:
            conn.close()


def extract_lecture_images(lecture_id):
    """Extract images for a lecture. Returns (count, source_type)."""
    conn = get_conn()
    row = conn.execute(
        "SELECT filename, source FROM lectures WHERE id=?", (lecture_id,)
    ).fetchone()
    if not row:
        conn.close()
        raise ValueError(f"No lecture {lecture_id}")
    _clear_slide_images(conn, lecture_id)
    conn.commit()

    lectures_dir = os.path.join(os.path.dirname(DB_PATH), "..", "lectures")
    filepath = os.path.join(lectures_dir, row["filename"])
    if row["source"] == "pdf":
        count = extract_pdf_images(lecture_id, filepath, conn=conn)
        kind = "page"
    else:
        count = extract_pptx_images(lecture_id, filepath, conn=conn)
        kind = "embedded"
    conn.commit()
    conn.close()
    return count, kind


def images_for_lecture(lecture_id):
    """Images grouped by slide, for the summary gallery."""
    conn = get_conn()
    rows = conn.execute(
        "SELECT si.*, s.slide_num FROM slide_images si "
        "JOIN slides s ON s.id = si.slide_id "
        "WHERE s.lecture_id=? ORDER BY s.slide_num, si.seq",
        (lecture_id,),
    ).fetchall()
    conn.close()
    grouped = {}
    for r in rows:
        grouped.setdefault(r["slide_num"], []).append(
            {"path": "/images/" + r["path"], "kind": r["kind"]}
        )
    return [{"slide_num": k, "images": v} for k, v in sorted(grouped.items())]


def images_for_slides(slide_ids):
    """Image URLs for a list of slide ids (for question source material)."""
    if not slide_ids:
        return []
    conn = get_conn()
    ph = ",".join("?" * len(slide_ids))
    rows = conn.execute(
        f"SELECT path FROM slide_images WHERE slide_id IN ({ph}) ORDER BY seq",
        tuple(slide_ids),
    ).fetchall()
    conn.close()
    return ["/images/" + r["path"] for r in rows]
