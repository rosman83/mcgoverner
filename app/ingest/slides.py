import os
import re
import fitz
from pptx import Presentation

from app.db import get_conn


def extract_pdf(path):
    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        text = re.sub(r"\s+", " ", text).strip()
        pages.append((i + 1, text))
    doc.close()
    return pages


def extract_pptx(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(cells))
        text = " ".join(parts)
        text = re.sub(r"\s+", " ", text).strip()
        slides.append((i + 1, text))
    return slides


def slugify(name):
    name = re.sub(r"[^\w\s-]", "", name).replace(" ", "_")
    return name.strip("_")


def import_lecture(path):
    ext = os.path.splitext(path)[1].lower()
    base = os.path.basename(path)
    title = os.path.splitext(base)[0]

    # Slow extraction first, before touching the DB, so no lock is held during it.
    slides = extract_pdf(path) if ext == ".pdf" else extract_pptx(path)

    conn = get_conn()
    cur = conn.execute(
        "INSERT INTO lectures(title, filename, source) VALUES(?,?,?)",
        (title, base, "pdf" if ext == ".pdf" else "pptx"),
    )
    lecture_id = cur.lastrowid

    word_count = 0
    for num, text in slides:
        conn.execute(
            "INSERT INTO slides(lecture_id, slide_num, text) VALUES(?,?,?)",
            (lecture_id, num, text),
        )
        word_count += len(text.split())
    conn.execute(
        "UPDATE lectures SET slide_count=?, word_count=? WHERE id=?",
        (len(slides), word_count, lecture_id),
    )
    conn.commit()
    conn.close()

    # Extract images after slides exist (best-effort; never blocks import)
    try:
        from app.ingest.images import extract_lecture_images
        extract_lecture_images(lecture_id)
        from app.ingest.ocr import run_ocr_for_lecture
        run_ocr_for_lecture(lecture_id)
    except Exception as e:
        print(f"image/ocr extraction for lecture {lecture_id} failed: {e}")

    return lecture_id


def list_lectures():
    conn = get_conn()
    rows = conn.execute(
        "SELECT * FROM lectures ORDER BY created_at DESC"
    ).fetchall()
    conn.close()
    return rows
